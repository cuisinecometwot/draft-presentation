#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import os
import fitz  # PyMuPDF
from pptx import Presentation
from PIL import Image
import io
from PyQt6.QtGui import QPixmap
from PyQt6.QtCore import QSize, Qt

class PresentationViewer:
    def __init__(self, file_path):
        self.file_path = file_path
        self.file_type = self._get_file_type(file_path)
        self.current_slide = 0
        self.slides = []
        self.zoom_factor = 1.0
        self.min_zoom = 0.1  # Zoom out tối thiểu (10%)
        self.max_zoom = 5.0  # Zoom in tối đa (500%)
        
        self._load_presentation()
        
    def _get_file_type(self, file_path):
        """Xác định loại file dựa trên phần mở rộng"""
        ext = os.path.splitext(file_path)[1].lower()
        if ext == '.pdf':
            return 'pdf'
        elif ext in ['.pptx', '.ppt']:
            return 'powerpoint'
        else:
            raise ValueError(f"Không hỗ trợ định dạng file: {ext}")
            
    def _load_presentation(self):
        """Tải presentation vào bộ nhớ"""
        if self.file_type == 'pdf':
            self._load_pdf()
        elif self.file_type == 'powerpoint':
            self._load_powerpoint()
            
    def _load_pdf(self):
        """Tải file PDF"""
        try:
            doc = fitz.open(self.file_path)
            self.slides = []
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                
                # Render page với độ phân giải cao
                mat = fitz.Matrix(2.0, 2.0)  # Scale 2x để có chất lượng tốt
                pix = page.get_pixmap(matrix=mat)
                
                # Chuyển đổi sang PIL Image
                img_data = pix.tobytes("ppm")
                img = Image.open(io.BytesIO(img_data))
                
                # Chuyển đổi sang QPixmap
                qpixmap = self._pil_to_qpixmap(img)
                self.slides.append(qpixmap)
                
            doc.close()
            
        except Exception as e:
            raise Exception(f"Lỗi khi tải PDF: {str(e)}")
            
    def _load_powerpoint(self):
        """Tải file PowerPoint"""
        try:
            prs = Presentation(self.file_path)
            self.slides = []
            
            for slide in prs.slides:
                # Tạo ảnh từ slide
                slide_image = self._render_slide_to_image(slide)
                self.slides.append(slide_image)
                
        except Exception as e:
            raise Exception(f"Lỗi khi tải PowerPoint: {str(e)}")
            
    def _render_slide_to_image(self, slide):
        """Render slide PowerPoint thành ảnh"""
        # Đây là một implementation đơn giản
        # Trong thực tế, bạn có thể cần sử dụng thư viện khác để render tốt hơn
        
        # Tạo một ảnh trống với kích thước chuẩn
        width, height = 1920, 1080  # 16:9 aspect ratio
        
        # Tạo QPixmap trống
        pixmap = QPixmap(width, height)
        pixmap.fill()  # Fill với màu trắng
        
        # Trong implementation thực tế, bạn sẽ render nội dung slide lên đây
        # Sử dụng thư viện như python-pptx để extract text và hình ảnh
        
        return pixmap
        
    def _pil_to_qpixmap(self, pil_image):
        """Chuyển đổi PIL Image sang QPixmap"""
        # Chuyển đổi PIL Image sang bytes
        buffer = io.BytesIO()
        pil_image.save(buffer, format='PNG')
        buffer.seek(0)
        
        # Tạo QPixmap từ bytes
        pixmap = QPixmap()
        pixmap.loadFromData(buffer.getvalue())
        
        return pixmap
        
    def get_current_slide(self, target_size=None):
        """Lấy slide hiện tại với zoom và fit size"""
        if not self.slides or not (0 <= self.current_slide < len(self.slides)):
            return None
            
        original_pixmap = self.slides[self.current_slide]
        
        # Nếu có target_size, tính toán fit size
        if target_size:
            return self._fit_slide_to_size(original_pixmap, target_size)
        
        # Nếu không có target_size, trả về slide gốc
        return original_pixmap
        
    def _fit_slide_to_size(self, pixmap, target_size):
        """Fit slide vào kích thước target với tỷ lệ khung hình"""
        if not pixmap or not target_size:
            return pixmap
            
        slide_width = pixmap.width()
        slide_height = pixmap.height()
        target_width = target_size.width()
        target_height = target_size.height()
        
        # Tính tỷ lệ để fit width hoặc height
        width_ratio = target_width / slide_width
        height_ratio = target_height / slide_height
        
        # Chọn tỷ lệ nhỏ hơn để fit toàn bộ slide
        scale_factor = min(width_ratio, height_ratio)
        
        # Áp dụng zoom factor
        final_scale = scale_factor * self.zoom_factor
        
        # Giới hạn zoom
        final_scale = max(self.min_zoom, min(self.max_zoom, final_scale))
        
        # Tính kích thước cuối cùng
        final_width = int(slide_width * final_scale)
        final_height = int(slide_height * final_scale)
        
        # Scale pixmap
        return pixmap.scaled(final_width, final_height, 
                           Qt.AspectRatioMode.KeepAspectRatio, 
                           Qt.TransformationMode.SmoothTransformation)
        
    def zoom_in(self, factor=1.2):
        """Zoom in slide"""
        self.zoom_factor = min(self.max_zoom, self.zoom_factor * factor)
        
    def zoom_out(self, factor=1.2):
        """Zoom out slide"""
        self.zoom_factor = max(self.min_zoom, self.zoom_factor / factor)
        
    def reset_zoom(self):
        """Reset zoom về mặc định"""
        self.zoom_factor = 1.0
        
    def get_zoom_factor(self):
        """Lấy zoom factor hiện tại"""
        return self.zoom_factor
        
    def get_current_slide_number(self):
        """Lấy số thứ tự slide hiện tại (0-based)"""
        return self.current_slide
        
    def get_total_slides(self):
        """Lấy tổng số slide"""
        return len(self.slides)
        
    def next_slide(self):
        """Chuyển đến slide tiếp theo"""
        if self.current_slide < len(self.slides) - 1:
            self.current_slide += 1
            
    def previous_slide(self):
        """Chuyển đến slide trước đó"""
        if self.current_slide > 0:
            self.current_slide -= 1
            
    def go_to_slide(self, slide_number):
        """Chuyển đến slide cụ thể"""
        if 0 <= slide_number < len(self.slides):
            self.current_slide = slide_number
            
    def get_slide_info(self):
        """Lấy thông tin về slide hiện tại"""
        if self.slides:
            current_pixmap = self.slides[self.current_slide]
            return {
                'current_slide': self.current_slide + 1,
                'total_slides': len(self.slides),
                'width': current_pixmap.width(),
                'height': current_pixmap.height(),
                'zoom_factor': self.zoom_factor
            }
        return None 